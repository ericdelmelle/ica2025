from pptx import Presentation
from pptx.util import Inches
import os

# Create a new presentation
prs = Presentation()

# Path to your folder with PNG files
img_folder = "/Users/delmelle/UNC Charlotte Dropbox/Eric Delmelle/XiaPaper/ica-2025/PPT"

# Get all png files with full path
png_files = [os.path.join(img_folder, f) for f in os.listdir(img_folder) if f.lower().endswith(".png")]

# Sort by modification time
png_files.sort(key=os.path.getmtime)

# Loop through sorted files
for img_path in png_files:
    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]  # 6 = blank
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add image full-slide
    slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                             width=prs.slide_width, height=prs.slide_height)

# Save the presentation
prs.save("output.pptx")
